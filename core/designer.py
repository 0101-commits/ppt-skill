# -*- coding: utf-8 -*-
"""HCG slide engine v2.0 — PPTX renderer.

Implements HCG-Slide-Design-System v1.0 (the absolute standard):
  - 16:9 fixed canvas, 1280x720 px  (= 13.333" x 7.5" @ 96dpi)
  - Pretendard ONLY (run XML a:latin + a:ea)
  - single main BLUE palette (#356CB5); maroon = "HCG" wordmark only
  - rigid 7-element skeleton inherited by every content slide:
      kicker / chapter-indicator / slide-title / title-rule /
      content-area / source-note / page-number
  - conclusion-sentence titles, charts-first, source always, >=85% density

Coordinates are px (top-left origin). PX(v) = Inches(v/96). Font size in points
= px * 0.75 (96dpi). The deck is built on a BLANK 16:9 canvas — the old 4:3
template-inheritance model is gone (legacy engine frozen at
history/designer_legacy_4x3.py).

The companion HTML renderer is core/html_renderer.py.
"""

import os
import sys

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Design tokens — single source core/tokens.py (shared with html_renderer).
# apply_design_tokens mutates the shared PALETTE dict in place.
from core.tokens import (PALETTE, TYPO, SAFE, CONTENT, GUTTER, BLOCK_GAP,
                         ROUNDED, FONT, SERIES_PALETTE, ROMAN,
                         apply_design_tokens)


def apply_theme(name=None):
    """Compatibility shim. v2.0 has a single fixed house style, so this is a
    no-op that returns the active palette (client brand is NOT a palette swap)."""
    return PALETTE


# ════════════════════════════════════════════════════════════
# Unit / color / text primitives
# ════════════════════════════════════════════════════════════

def PX(v):
    """px -> EMU length (Inches at 96 dpi)."""
    return Inches(v / 96.0)


def PTS(px):
    """css px -> points (96dpi): pt = px * 0.75."""
    return Pt(px * 0.75)


def _rgb(name_or_hex):
    s = PALETTE.get(name_or_hex, name_or_hex)
    s = str(s).lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _hex6(name_or_hex):
    return str(PALETTE.get(name_or_hex, name_or_hex)).lstrip("#").upper()


def _set_font_xml(run, name=FONT):
    """Force Pretendard on latin + east-asian glyphs (prevents font fallback)."""
    rPr = run._r.get_or_add_rPr()
    for tag in (qn("a:latin"), qn("a:ea"), qn("a:cs")):
        el = rPr.find(tag)
        if el is None:
            el = etree.SubElement(rPr, tag)
        el.set("typeface", name)


def _segments(text):
    """Parse '**emphasis**' markup -> [(chunk, is_emph), ...]."""
    out = []
    for i, part in enumerate(str(text).split("**")):
        if part:
            out.append((part, i % 2 == 1))
    return out or [("", False)]


def _line_spacing(p, lh):
    pPr = p._p.get_or_add_pPr()
    for old in pPr.findall(qn("a:lnSpc")):
        pPr.remove(old)
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    pct = etree.SubElement(lnSpc, qn("a:spcPct"))
    pct.set("val", str(int(lh * 100000)))
    pPr.insert(0, lnSpc)


def _run(p, text, role, color, bold=None, italic=False):
    px, wt, _lh = TYPO[role]
    r = p.add_run()
    r.text = text
    r.font.size = PTS(px)
    r.font.bold = (wt >= 700) if bold is None else bold
    r.font.italic = italic
    r.font.color.rgb = _rgb(color)
    _set_font_xml(r)
    return r


def _textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(PX(x), PX(y), PX(w), PX(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = anchor
    return tb, tf


def _write(slide, x, y, w, h, role, text, color="ink", align=PP_ALIGN.LEFT,
           anchor=MSO_ANCHOR.TOP, emph_color="primary", italic=False):
    """One-paragraph textbox honoring a typography role; supports **emph** runs."""
    _tb, tf = _textbox(slide, x, y, w, h, anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    _line_spacing(p, TYPO[role][2])
    for chunk, emph in _segments(text):
        _run(p, chunk, role, emph_color if emph else color, italic=italic)
    return tf


def _box(slide, x, y, w, h, fill=None, line=None, line_px=1.0, rounded=None):
    shp = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shp, PX(x), PX(y), PX(w), PX(h))
    if rounded:
        try:
            s.adjustments[0] = max(0.0, min(0.5, rounded / float(min(w, h))))
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = _rgb(line)
        s.line.width = PTS(line_px)
    s.shadow.inherit = False
    return s


def _shape_text(shape, role, text, color="ink", align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE, emph_color="primary"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    _line_spacing(p, TYPO[role][2])
    for chunk, emph in _segments(text):
        _run(p, chunk, role, emph_color if emph else color)
    return tf


def _hline(slide, x1, y, x2, color="divider", px=0.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, PX(x1), PX(y), PX(x2), PX(y))
    c.line.color.rgb = _rgb(color)
    c.line.width = PTS(px)
    c.shadow.inherit = False
    return c


# ════════════════════════════════════════════════════════════
# Rigid skeleton (inherited by every content slide)
# ════════════════════════════════════════════════════════════

def draw_skeleton(slide, title=None, kicker=None, chapter=None, source=None,
                  page=None):
    """Lay the 7 fixed elements at their absolute px coordinates."""
    if kicker:
        # whole kicker is a topic label: emphasize keyword(s) primary, rest ink
        has_emph = "**" in kicker
        _write(slide, 56, 28, 760, 28, "kicker", kicker,
               color="ink" if has_emph else "primary",
               anchor=MSO_ANCHOR.MIDDLE)
    if chapter:
        cw = 150
        cb = _box(slide, 1224 - cw, 28, cw, 26, fill="surface",
                  line="divider", line_px=1.0, rounded=ROUNDED["md"])
        _shape_text(cb, "chapter-indicator", chapter, color="primary-steel")
    if title is not None:
        _write(slide, 56, 64, 1168, 64, "slide-title", title, color="ink",
               anchor=MSO_ANCHOR.MIDDLE, emph_color="primary")
        _hline(slide, 56, 136, 1224, color="divider", px=1.0)
    if source:
        txt = source if str(source).strip().startswith("※") else "※ Source : " + source
        _write(slide, 56, 678, 980, 16, "source-note", txt, color="text-muted",
               anchor=MSO_ANCHOR.MIDDLE)
    if page:
        _tb, tf = _textbox(slide, 1024, 676, 200, 18, anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        _line_spacing(p, TYPO["page-number"][2])
        _run(p, f"{page}   ", "page-number", "text-muted")
        _run(p, "HCG", "page-number", "hcg-maroon", bold=True)


# ════════════════════════════════════════════════════════════
# Content builders  (draw inside content-area 56,152,1168,512)
# ════════════════════════════════════════════════════════════

def _grid_cols(n, x=CONTENT["x"], w=CONTENT["w"], gap=GUTTER):
    n = max(1, n)
    cw = (w - gap * (n - 1)) / n
    return [(x + i * (cw + gap), cw) for i in range(n)], cw


def _bullets_into(slide, items, x, y, w, h, role="body"):
    """■ bullet list with optional one level of sub-items."""
    _tb, tf = _textbox(slide, x, y, w, h)
    first = True
    for it in items:
        text = it if isinstance(it, str) else it.get("text", "")
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        _line_spacing(p, 1.4)
        p.space_after = Pt(4)
        _run(p, "■ ", role, "primary", bold=True)
        _run(p, text, role, "ink")
        if isinstance(it, dict):
            for sub in it.get("sub", []) or []:
                sp = tf.add_paragraph()
                sp.alignment = PP_ALIGN.LEFT
                sp.level = 1
                _line_spacing(sp, 1.35)
                _run(sp, "– ", "body-sm", "text-muted")
                _run(sp, sub, "body-sm", "ink-slate")
    return tf


def build_bullets(slide, s):
    _bullets_into(slide, s["items"], CONTENT["x"], CONTENT["y"] + 8,
                  CONTENT["w"], CONTENT["h"] - 8)


def build_cards(slide, s):
    items = s["items"]
    n = len(items)
    per_row = n if n <= 4 else (n + 1) // 2
    rows = 1 if n <= 4 else 2
    cells, cw = _grid_cols(per_row)
    row_h = (CONTENT["h"] - BLOCK_GAP * (rows - 1)) / rows
    for i, it in enumerate(items):
        r, c = divmod(i, per_row)
        cx, _ = cells[c] if c < len(cells) else cells[-1]
        cy = CONTENT["y"] + r * (row_h + BLOCK_GAP)
        card = _box(slide, cx, cy, cw, row_h, fill="surface", line="divider",
                    line_px=1.0, rounded=ROUNDED["card"])
        header = it.get("header", "") if isinstance(it, dict) else ""
        hb_h = 34
        if header:
            _box(slide, cx, cy, cw, hb_h, fill="divider-soft", line=None,
                 rounded=ROUNDED["sm"])
            _shape_text(_box(slide, cx, cy, cw, hb_h, fill=None),
                        "block-header", header, color="ink")
            body_y, body_h = cy + hb_h + 6, row_h - hb_h - 14
        else:
            body_y, body_h = cy + 10, row_h - 20
        if isinstance(it, dict) and it.get("bullets"):
            _bullets_into(slide, it["bullets"], cx + 12, body_y, cw - 24, body_h,
                          role="body-sm")
        else:
            body = it.get("body", "") if isinstance(it, dict) else str(it)
            _write(slide, cx + 12, body_y, cw - 24, body_h, "body-sm", body,
                   color="ink-slate")


def build_columns(slide, s):
    cols = s["columns"]
    cells, cw = _grid_cols(len(cols))
    for (cx, _), col in zip(cells, cols):
        dark = col.get("dark")
        hb = _box(slide, cx, CONTENT["y"], cw, 38,
                  fill="primary-deep" if dark else "primary-soft",
                  rounded=ROUNDED["sm"])
        _shape_text(hb, "section-header", col.get("header", ""),
                    color="surface" if dark else "ink")
        body_y = CONTENT["y"] + 38 + BLOCK_GAP
        body_h = CONTENT["h"] - 38 - BLOCK_GAP
        if col.get("bullets"):
            _bullets_into(slide, col["bullets"], cx + 6, body_y, cw - 12, body_h)
        else:
            _box(slide, cx, body_y, cw, body_h, fill="surface", line="divider",
                 rounded=ROUNDED["card"])
            _write(slide, cx + 14, body_y + 12, cw - 28, body_h - 24, "body",
                   col.get("body", ""), color="ink-slate")


def build_compare(slide, s):
    left, right = s["left"], s["right"]
    (lx, _), (rx, _) = _grid_cols(2)[0]
    cw = _grid_cols(2)[1]
    headers = [(lx, s.get("left_header", "As-Is"), False),
               (rx, s.get("right_header", "To-Be"), True)]
    has_take = bool(s.get("takeaway"))
    list_top = CONTENT["y"] + 38 + BLOCK_GAP
    list_h = CONTENT["h"] - 38 - BLOCK_GAP - (40 if has_take else 0)
    for x, htext, dark in headers:
        hb = _box(slide, x, CONTENT["y"], cw, 38,
                  fill="primary-deep" if dark else "primary-soft",
                  rounded=ROUNDED["sm"])
        _shape_text(hb, "section-header", htext,
                    color="surface" if dark else "ink")
    for x, rows, dark in ((lx, left, False), (rx, right, True)):
        n = max(1, len(rows))
        rh = (list_h - BLOCK_GAP * (n - 1)) / n
        for i, row in enumerate(rows):
            ry = list_top + i * (rh + BLOCK_GAP)
            cell = _box(slide, x, ry, cw, rh,
                        fill="tint-blue-4" if dark else "surface",
                        line="tint-blue-3" if dark else "divider",
                        rounded=ROUNDED["card"])
            _shape_text(cell, "body-sm", row, color="ink",
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    if has_take:
        ty = CONTENT["y"] + CONTENT["h"] - 32
        chip = _box(slide, CONTENT["x"] + 200, ty, CONTENT["w"] - 400, 30,
                    fill="accent-coral", rounded=ROUNDED["chip"])
        _shape_text(chip, "body-sm-strong", s["takeaway"], color="surface")


def build_kpi(slide, s):
    items = s["items"]
    cells, cw = _grid_cols(len(items))
    ch = 200
    cy = CONTENT["y"] + (CONTENT["h"] - ch) / 2
    for (cx, _), it in zip(cells, items):
        _box(slide, cx, cy, cw, ch, fill="tint-blue-4", line="tint-blue-3",
             rounded=ROUNDED["card"])
        _tb, tf = _textbox(slide, cx + 10, cy + 30, cw - 20, ch - 60,
                           anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _line_spacing(p, 1.05)
        _run(p, str(it.get("number", "")), "kpi-number", "primary")
        if it.get("unit"):
            _run(p, " " + it["unit"], "kpi-unit", "text-muted")
        lp = tf.add_paragraph()
        lp.alignment = PP_ALIGN.CENTER
        _line_spacing(lp, 1.3)
        _run(lp, it.get("label", ""), "caption", "ink-slate")


def build_process(slide, s):
    steps = s["steps"]
    n = len(steps)
    arrow_w = 30
    cells, cw = _grid_cols(n, gap=arrow_w + 8)
    node_h = 70
    ny = CONTENT["y"] + 60
    for i, ((cx, _), st) in enumerate(zip(cells, steps)):
        node = _box(slide, cx, ny, cw, node_h, fill="tint-blue-3",
                    rounded=ROUNDED["card"])
        _shape_text(node, "body-sm-strong", st.get("label", ""), color="ink")
        if st.get("desc"):
            _write(slide, cx, ny + node_h + 10, cw, 140, "body-sm",
                   st["desc"], color="ink-slate", align=PP_ALIGN.CENTER)
        if i < n - 1:
            ax = cx + cw + 4
            ar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                                        PX(ax), PX(ny + node_h / 2 - 9),
                                        PX(arrow_w), PX(18))
            ar.fill.solid()
            ar.fill.fore_color.rgb = _rgb("text-muted")
            ar.line.fill.background()
            ar.shadow.inherit = False


def build_matrix(slide, s):
    cols = s["columns"]
    rows = s["rows"]
    has_take = bool(s.get("takeaway"))
    ncol = len(cols)
    table_h = CONTENT["h"] - (44 if has_take else 0)
    header_h = 40
    body_h = table_h - header_h
    rh = body_h / max(1, len(rows))
    col_w = CONTENT["w"] / ncol
    for j, ch in enumerate(cols):
        hb = _box(slide, CONTENT["x"] + j * col_w, CONTENT["y"], col_w, header_h,
                  fill="tint-blue-1", line="surface", line_px=1.0)
        _shape_text(hb, "body-sm-strong", ch, color="ink")
    for i, row in enumerate(rows):
        ry = CONTENT["y"] + header_h + i * rh
        cells = row.get("cells", []) if isinstance(row, dict) else list(row)
        group = row.get("group", "") if isinstance(row, dict) else ""
        values = [group] + list(cells) if group else list(cells)
        zebra = "surface" if i % 2 == 0 else "surface-gray"
        for j in range(ncol):
            cx = CONTENT["x"] + j * col_w
            is_group = (j == 0 and group)
            _box(slide, cx, ry, col_w, rh,
                 fill="tint-blue-4" if is_group else zebra,
                 line="divider", line_px=0.75)
            val = values[j] if j < len(values) else ""
            _shape_text(_box(slide, cx, ry, col_w, rh, fill=None),
                        "body-sm-strong" if is_group else "body-sm", val,
                        color="primary" if is_group else "ink",
                        align=PP_ALIGN.CENTER if is_group else PP_ALIGN.LEFT)
    if has_take:
        ty = CONTENT["y"] + CONTENT["h"] - 34
        chip = _box(slide, CONTENT["x"] + 160, ty, CONTENT["w"] - 320, 30,
                    fill="accent-coral", rounded=ROUNDED["chip"])
        _shape_text(chip, "body-sm-strong", s["takeaway"], color="surface")


def build_table(slide, s):
    headers, rows = s["headers"], s["rows"]
    emph = s.get("emphasis_row")
    ncol = len(headers)
    nrow = len(rows) + 1
    h = min(CONTENT["h"], 44 * nrow + 8)
    gf = slide.shapes.add_table(nrow, ncol, PX(CONTENT["x"]), PX(CONTENT["y"]),
                                PX(CONTENT["w"]), PX(h))
    tbl = gf.table
    data = [headers] + rows
    for r in range(nrow):
        for c in range(ncol):
            cell = tbl.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text = str(data[r][c]) if c < len(data[r]) else ""
            is_head = (r == 0)
            is_emph = (emph is not None and r == emph + 1)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(
                "tint-blue-1" if is_head else "tint-blue-2" if is_emph
                else ("surface" if r % 2 else "surface-gray"))
            role = "body-sm-strong" if (is_head or is_emph) else "body-sm"
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if (is_head or c > 0) else PP_ALIGN.LEFT
                _line_spacing(p, TYPO[role][2])
                for run in p.runs:
                    run.font.size = PTS(TYPO[role][0])
                    run.font.bold = is_head or is_emph
                    run.font.color.rgb = _rgb("ink")
                    _set_font_xml(run)
    return gf


# ── charts (native python-pptx) ──────────────────────────────
from pptx.chart.data import CategoryChartData            # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # noqa: E402

_CHART_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked-bar": XL_CHART_TYPE.BAR_STACKED,
    "stacked-column": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}


def build_chart(slide, s):
    spec = s["chart"]
    ctype = _CHART_MAP.get(spec.get("type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    cdata = CategoryChartData()
    cdata.categories = spec.get("labels", [])
    series = spec.get("series", [])
    for ser in series:
        cdata.add_series(ser.get("name", ""), ser.get("data", []))
    gf = slide.shapes.add_chart(ctype, PX(CONTENT["x"]), PX(CONTENT["y"]),
                                PX(CONTENT["w"]), PX(CONTENT["h"]), cdata)
    chart = gf.chart
    try:
        chart.font.name = FONT
        chart.font.size = PTS(TYPO["chart-axis"][0])
    except Exception:
        pass
    is_single = spec.get("type") in ("pie", "donut", "doughnut")
    chart.has_legend = (len(series) > 1) or is_single
    if chart.has_legend:
        try:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        except Exception:
            pass
    # series colors per the fixed palette
    try:
        plots = chart.plots
        if is_single:
            pts = plots[0].series[0].points
            for i, pt in enumerate(pts):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = _rgb(SERIES_PALETTE[i % len(SERIES_PALETTE)])
        else:
            for i, ser in enumerate(chart.series):
                col = _rgb(SERIES_PALETTE[i % len(SERIES_PALETTE)])
                if spec.get("type") in ("line",):
                    ser.format.line.color.rgb = col
                    ser.format.line.width = Pt(2.25)
                else:
                    ser.format.fill.solid()
                    ser.format.fill.fore_color.rgb = col
    except Exception:
        pass
    return gf


# ════════════════════════════════════════════════════════════
# Structural slides
# ════════════════════════════════════════════════════════════

def build_cover(slide, s):
    _box(slide, 0, 0, 1280, 300, fill="primary")
    _box(slide, 0, 296, 1280, 6, fill="primary-deep")
    if s.get("confidential", True):
        _write(slide, 56, 20, 400, 22, "caption", "Strictly Confidential",
               color="surface", italic=True)
    _write(slide, 160, 392, 960, 56, "cover-title", s["title"], color="ink",
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, emph_color="primary")
    if s.get("subtitle"):
        _write(slide, 160, 470, 960, 30, "cover-subtitle", s["subtitle"],
               color="primary-steel", align=PP_ALIGN.CENTER)
    if s.get("date"):
        _write(slide, 160, 540, 960, 24, "cover-date", s["date"],
               color="text-muted", align=PP_ALIGN.CENTER)
    _write(slide, 560, 612, 160, 48, "section-header", "HCG",
           color="hcg-maroon", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    legal = s.get("legal") or ("본 제안서는 ㈜휴먼컨설팅그룹의 편집저작물로서 "
                               "관련 법령에 의해 보호됩니다.")
    _write(slide, 160, 676, 960, 28, "source-note", legal, color="text-muted",
           align=PP_ALIGN.CENTER)


def build_toc(slide, s):
    items = s["items"]
    current = s.get("current")
    _write(slide, 56, 64, 1168, 64, "slide-title", s.get("title", "목차"),
           color="ink", anchor=MSO_ANCHOR.MIDDLE)
    _hline(slide, 56, 136, 1224, color="divider", px=1.0)
    n = len(items)
    row_h = min(86, (CONTENT["h"]) / n)
    top = CONTENT["y"] + (CONTENT["h"] - row_h * n) / 2
    for i, name in enumerate(items):
        active = (current == i)
        y = top + i * row_h
        chip = _box(slide, 320, y + (row_h - 52) / 2, 64, 52,
                    fill="primary" if active else "tint-blue-2",
                    rounded=ROUNDED["sm"])
        _shape_text(chip, "section-header",
                    ROMAN[i] if i < len(ROMAN) else str(i + 1),
                    color="surface")
        _write(slide, 408, y, 560, row_h, "section-header", name,
               color="primary" if active else "ink",
               anchor=MSO_ANCHOR.MIDDLE)
        if active:
            _box(slide, 396, y + 6, 564, row_h - 12, fill=None,
                 line="primary", line_px=1.5, rounded=ROUNDED["sm"])


def build_section(slide, s):
    _box(slide, 0, 240, 1280, 240, fill="tint-blue-4")
    _write(slide, 96, 250, 300, 120, "kpi-number", s.get("num", ""),
           color="primary", anchor=MSO_ANCHOR.MIDDLE)
    _write(slide, 300, 300, 884, 56, "cover-title", s["title"], color="ink",
           anchor=MSO_ANCHOR.MIDDLE, emph_color="primary")
    if s.get("sub"):
        _write(slide, 300, 372, 884, 60, "cover-subtitle", s["sub"],
               color="ink-slate")


def build_end(slide, s):
    _box(slide, 0, 0, 1280, 720, fill="primary")
    _write(slide, 160, 300, 960, 80, "cover-title", s.get("message", "Thank You"),
           color="surface", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _write(slide, 160, 410, 960, 30, "cover-subtitle", "㈜휴먼컨설팅그룹",
           color="tint-blue-2", align=PP_ALIGN.CENTER)


_CONTENT_BUILDERS = {
    "bullets": build_bullets, "cards": build_cards, "columns": build_columns,
    "compare": build_compare, "kpi": build_kpi, "process": build_process,
    "matrix": build_matrix, "chart": build_chart, "table": build_table,
}


# ════════════════════════════════════════════════════════════
# Engine
# ════════════════════════════════════════════════════════════

class DeckEngine:
    CANVAS_W, CANVAS_H = 1280, 720

    def __init__(self, out=None, design_tokens=None, color_overrides=None,
                 **_ignored):
        # _ignored absorbs legacy kwargs (template/theme) — v2.0 builds a blank
        # 16:9 canvas and uses a single fixed house style.
        if design_tokens or color_overrides:
            apply_design_tokens(design_tokens or {}, color_overrides)
        self.out = out
        self.prs = Presentation()
        self.prs.slide_width = PX(self.CANVAS_W)
        self.prs.slide_height = PX(self.CANVAS_H)
        self.blank = self.prs.slide_layouts[6]   # blank layout in default master
        self.page = 0
        self.count = 0

    def _new(self):
        return self.prs.slides.add_slide(self.blank)

    def render(self, spec):
        for s in spec.get("slides", []):
            t = s.get("type")
            slide = self._new()
            self.count += 1
            if t == "cover":
                build_cover(slide, s)
            elif t == "section":
                build_section(slide, s)
            elif t == "end":
                build_end(slide, s)
            elif t == "toc":
                build_toc(slide, s)
                self.page += 1
                draw_skeleton(slide, source=s.get("source"), page=self.page)
            elif t in _CONTENT_BUILDERS:
                self.page += 1
                draw_skeleton(slide, title=s.get("title"), kicker=s.get("kicker"),
                              chapter=s.get("chapter"), source=s.get("source"),
                              page=self.page)
                _CONTENT_BUILDERS[t](slide, s)
            else:
                print("  [skip] unknown slide type:", t)
                self.count -= 1
                continue
            print(f"  S{self.count:02d} {s.get('label', t)}")
        return self

    def save(self, out=None):
        path = out or self.out
        self.prs.save(path)
        print(f"[saved pptx] {path}")
        return path


Designer = DeckEngine   # public alias for the framework
