# -*- coding: utf-8 -*-
"""Phase 1: Reverse-engineer the human final proposal pptx.
Extracts slide layouts, shape trees, coordinates (EMU->cm/inch), fonts, colors,
and detects grid / container (box-in-box) structure patterns.
"""
import sys, json, collections
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

FINAL = "C:/Users/cgpar/OneDrive - 휴먼컨설팅그룹/09 Admin/09 etc/other/Claude/기아 제안서/기아_중장기 보상체계 개선 추진_제안서_250912_v1.1_final.pptx"

EMU_PER_CM = 360000.0
EMU_PER_INCH = 914400.0


def emu_cm(v):
    if v is None:
        return None
    return round(v / EMU_PER_CM, 2)


def color_hex(color):
    try:
        if color and color.type is not None:
            return str(color.rgb)
    except Exception:
        return None
    return None


def font_info(font):
    d = {}
    try:
        d["name"] = font.name
        d["size"] = font.size.pt if font.size is not None else None
        d["bold"] = font.bold
        d["italic"] = font.italic
        d["color"] = color_hex(font.color)
    except Exception as e:
        d["err"] = str(e)
    return d


def shape_fill(shape):
    try:
        fill = shape.fill
        ft = str(fill.type)
        hexc = None
        try:
            hexc = str(fill.fore_color.rgb)
        except Exception:
            hexc = None
        return {"type": ft, "color": hexc}
    except Exception:
        return None


def shape_line(shape):
    try:
        ln = shape.line
        c = None
        try:
            c = str(ln.color.rgb)
        except Exception:
            c = None
        w = None
        try:
            w = ln.width.pt if ln.width is not None else None
        except Exception:
            w = None
        return {"color": c, "width_pt": w}
    except Exception:
        return None


def parse_text(tf):
    paras = []
    for p in tf.paragraphs:
        runs = []
        for r in p.runs:
            runs.append({"text": r.text, "font": font_info(r.font)})
        paras.append({
            "level": p.level,
            "align": str(p.alignment),
            "text": p.text,
            "runs": runs,
        })
    return paras


def walk_shape(shape, depth=0):
    info = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "id": shape.shape_id,
        "x_cm": emu_cm(shape.left),
        "y_cm": emu_cm(shape.top),
        "w_cm": emu_cm(shape.width),
        "h_cm": emu_cm(shape.height),
        "depth": depth,
    }
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            info["auto_shape"] = str(shape.adjustments) if False else None
            try:
                info["adj"] = str(shape.auto_shape_type)
            except Exception:
                pass
    except Exception:
        pass
    info["fill"] = shape_fill(shape)
    info["line"] = shape_line(shape)
    if shape.has_text_frame and shape.text_frame.text.strip():
        info["text"] = parse_text(shape.text_frame)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        info["children"] = [walk_shape(c, depth + 1) for c in shape.shapes]
    return info


def detect_grid(shapes):
    """Detect horizontally-aligned shape rows (grid arrangement)."""
    rows = collections.defaultdict(list)
    for s in shapes:
        if s["y_cm"] is None or s["w_cm"] is None:
            continue
        key = round(s["y_cm"])  # bucket by y
        rows[key].append(s)
    grids = []
    for y, items in rows.items():
        if len(items) >= 3:
            xs = sorted(items, key=lambda z: z["x_cm"] or 0)
            grids.append({
                "y_cm": y,
                "count": len(items),
                "x_positions": [i["x_cm"] for i in xs],
                "widths": [i["w_cm"] for i in xs],
                "names": [i["name"] for i in xs],
            })
    return grids


def main():
    prs = Presentation(FINAL)
    out = {
        "slide_w_cm": emu_cm(prs.slide_width),
        "slide_h_cm": emu_cm(prs.slide_height),
        "n_slides": len(prs.slides),
        "slides": [],
    }
    print("=== FINAL PPTX ANALYSIS ===")
    print(f"slide size: {out['slide_w_cm']} x {out['slide_h_cm']} cm | slides: {out['n_slides']}")

    for idx, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        top_shapes = []
        for sh in slide.shapes:
            top_shapes.append(walk_shape(sh))
        grids = detect_grid(top_shapes)
        # collect first-line text to guess slide type
        title_text = ""
        for s in top_shapes:
            if "text" in s and s["text"]:
                title_text = s["text"][0]["text"][:40]
                break
        sl = {
            "idx": idx,
            "layout": layout_name,
            "n_shapes": len(top_shapes),
            "title_guess": title_text,
            "grids": grids,
            "shapes": top_shapes,
        }
        out["slides"].append(sl)
        gflag = f" GRID({len(grids)})" if grids else ""
        print(f"[{idx:02d}] layout='{layout_name}' shapes={len(top_shapes)}{gflag} :: {title_text}")

    with open(r"C:\Users\cgpar\ppt-skill\final_analysis.json", "w", encoding="utf-8") as f:
        json.dump(out, f, ensure_ascii=False, indent=1)
    print("\nwrote final_analysis.json")

    # Identify TOC slide candidates (contains 목차/Contents/Agenda or many number markers)
    print("\n=== TOC CANDIDATES ===")
    for sl in out["slides"]:
        joined = " ".join(
            s.get("text", [{}])[0].get("text", "") for s in sl["shapes"] if s.get("text")
        )
        if any(k in joined for k in ["목차", "Contents", "CONTENTS", "Agenda", "AGENDA"]):
            print(f"  slide {sl['idx']}: {sl['title_guess']}  (shapes={sl['n_shapes']}, grids={len(sl['grids'])})")


if __name__ == "__main__":
    main()
