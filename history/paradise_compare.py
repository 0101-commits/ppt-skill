# -*- coding: utf-8 -*-
"""Phase 1/2: AI Draft vs 인간 rev_final 정밀 비교.
브랜드컬러/윤곽선/행간/그리드 좌표 오차 + 텍스트 오버플로 탐지.
"""
import sys, os, json, math, collections
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

FOLDER = r"C:\Users\cgpar\OneDrive - 휴먼컨설팅그룹\09 Admin\09 etc\other\Claude\파라다이스 제안서"
AI   = os.path.join(FOLDER, "HCG_파라다이스_인사제도_개선_제안서_AI_Draft.pptx")
HUMAN = os.path.join(FOLDER, "HCG_파라다이스_보상제도 컨설팅_제안서_rev_final.pptx")

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
EMU_CM = 360000.0
EMU_PT = 12700.0


def cm(v):
    return round(v / EMU_CM, 2) if v is not None else None


def pt(v):
    return round(v / EMU_PT, 1) if v is not None else None


def run_color(r):
    rPr = r._r.find(qn('a:rPr'))
    if rPr is None:
        return None
    sf = rPr.find(qn('a:solidFill'))
    if sf is None:
        return None
    c = sf.find(qn('a:srgbClr'))
    return c.get('val').upper() if c is not None else None


def shape_fill_hex(sh):
    try:
        if sh.fill.type is not None:
            return str(sh.fill.fore_color.rgb).upper()
    except Exception:
        return None
    return None


def stats_for(path, label):
    prs = Presentation(path)
    st = {
        "label": label,
        "slides": len(prs.slides),
        "w_cm": cm(prs.slide_width),
        "h_cm": cm(prs.slide_height),
        "fill_colors": collections.Counter(),
        "font_colors": collections.Counter(),
        "font_names": collections.Counter(),
        "font_sizes": collections.Counter(),
        "outline_runs": 0,        # <a:rPr><a:ln>
        "shadow": 0,              # <a:outerShdw>
        "transparency": 0,        # <a:alpha>
        "gradient": 0,            # <a:gradFill>
        "char_spacing": 0,        # spc attr
        "linespacing": collections.Counter(),
        "autofit": collections.Counter(),
        "shape_types": collections.Counter(),
        "overflow": [],           # (slide, name, need_pt, box_pt, text)
        "title_coords": [],       # (x,y,w,h) of title ph idx0
        "sub_coords": [],
    }

    def walk(sh, sidx):
        st["shape_types"][str(sh.shape_type)] += 1
        f = shape_fill_hex(sh)
        if f:
            st["fill_colors"][f] += 1
        el = sh._element
        # gradient
        st["gradient"] += len(el.findall('.//' + qn('a:gradFill')))
        st["shadow"] += len(el.findall('.//' + qn('a:outerShdw')))
        st["transparency"] += len(el.findall('.//' + qn('a:alpha')))
        # placeholder coords
        try:
            if sh.is_placeholder:
                pid = sh.placeholder_format.idx
                if pid == 0:
                    st["title_coords"].append((cm(sh.left), cm(sh.top), cm(sh.width), cm(sh.height)))
                elif pid in (1, 10):
                    st["sub_coords"].append((cm(sh.left), cm(sh.top), cm(sh.width), cm(sh.height)))
        except Exception:
            pass
        if sh.has_text_frame:
            tf = sh.text_frame
            body = tf._txBody
            bodyPr = body.find(qn('a:bodyPr'))
            if bodyPr is not None:
                if bodyPr.find(qn('a:normAutofit')) is not None:
                    st["autofit"]['normAutofit'] += 1
                elif bodyPr.find(qn('a:spAutoFit')) is not None:
                    st["autofit"]['spAutoFit'] += 1
                elif bodyPr.find(qn('a:noAutofit')) is not None:
                    st["autofit"]['noAutofit'] += 1
                else:
                    st["autofit"]['(none)'] += 1
            # text + font stats + overflow heuristic
            full = tf.text
            max_fs = 0
            nlines_text = 0
            for p in tf.paragraphs:
                # line spacing
                pPr = p._p.find(qn('a:pPr'))
                if pPr is not None:
                    ln = pPr.find(qn('a:lnSpc'))
                    if ln is not None:
                        pct = ln.find(qn('a:spcPct'))
                        if pct is not None:
                            st["linespacing"][int(int(pct.get('val')) / 1000)] += 1
                ptxt = p.text
                for r in p.runs:
                    rPr = r._r.find(qn('a:rPr'))
                    if rPr is not None:
                        if rPr.find(qn('a:ln')) is not None:
                            st["outline_runs"] += 1
                        if rPr.get('spc'):
                            st["char_spacing"] += 1
                        sz = rPr.get('sz')
                        if sz:
                            fs = int(sz) / 100.0
                            st["font_sizes"][fs] += 1
                            max_fs = max(max_fs, fs)
                        ln_el = rPr.find(qn('a:latin'))
                        if ln_el is not None and ln_el.get('typeface'):
                            st["font_names"][ln_el.get('typeface')] += 1
                    c = run_color(r)
                    if c:
                        st["font_colors"][c] += 1
                if ptxt.strip():
                    nlines_text += 1
            # overflow heuristic (Korean ~ 1 char width = font pt)
            if full.strip() and max_fs > 0 and sh.width and sh.height:
                box_w_pt = sh.width / EMU_PT
                box_h_pt = sh.height / EMU_PT
                # margins approx
                cpl = max(1, int(box_w_pt / (max_fs * 1.05)))
                # estimate wrapped lines per paragraph
                est_lines = 0
                for p in tf.paragraphs:
                    L = len(p.text)
                    est_lines += max(1, math.ceil(L / cpl)) if L else 1
                need_pt = est_lines * max_fs * 1.18  # line spacing ~118%
                if need_pt > box_h_pt * 1.08:
                    st["overflow"].append((sidx, sh.name, round(need_pt, 0),
                                           round(box_h_pt, 0), full.replace("\n", " ")[:40]))
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for c in sh.shapes:
                walk(c, sidx)

    for i, sl in enumerate(prs.slides):
        for sh in sl.shapes:
            walk(sh, i)
    return st


def main():
    for p, n in [(AI, "AI"), (HUMAN, "HUMAN")]:
        if not os.path.exists(p):
            print("MISSING:", p)
            return
    a = stats_for(AI, "AI_Draft")
    h = stats_for(HUMAN, "Human_rev_final")

    out = []
    out.append("# PARADISE  AI_Draft  vs  Human_rev_final  비교\n")
    out.append(f"{'항목':<22}{'AI':>16}{'HUMAN':>16}")
    out.append("-" * 54)

    def row(k, av, hv):
        out.append(f"{k:<22}{str(av):>16}{str(hv):>16}")

    row("슬라이드 수", a["slides"], h["slides"])
    row("크기(cm)", f"{a['w_cm']}x{a['h_cm']}", f"{h['w_cm']}x{h['h_cm']}")
    row("윤곽선 runs(a:ln)", a["outline_runs"], h["outline_runs"])
    row("그림자 outerShdw", a["shadow"], h["shadow"])
    row("투명도 alpha", a["transparency"], h["transparency"])
    row("그라디언트", a["gradient"], h["gradient"])
    row("자간 spc", a["char_spacing"], h["char_spacing"])
    ai_shapes = sum(a["shape_types"].values())
    hu_shapes = sum(h["shape_types"].values())
    row("총 도형", ai_shapes, hu_shapes)
    row("도형/장", round(ai_shapes / a["slides"], 1), round(hu_shapes / h["slides"], 1))

    out.append("\n## 브랜드 FILL 컬러 (top10)")
    out.append(f"AI   : {dict(a['fill_colors'].most_common(10))}")
    out.append(f"HUMAN: {dict(h['fill_colors'].most_common(10))}")
    out.append("\n## FONT 컬러 (top10)")
    out.append(f"AI   : {dict(a['font_colors'].most_common(10))}")
    out.append(f"HUMAN: {dict(h['font_colors'].most_common(10))}")
    out.append("\n## 폰트명 (top6)")
    out.append(f"AI   : {dict(a['font_names'].most_common(6))}")
    out.append(f"HUMAN: {dict(h['font_names'].most_common(6))}")
    out.append("\n## 폰트 크기 (top10)")
    out.append(f"AI   : {dict(a['font_sizes'].most_common(10))}")
    out.append(f"HUMAN: {dict(h['font_sizes'].most_common(10))}")
    out.append("\n## 행간 % 분포")
    out.append(f"AI   : {dict(a['linespacing'].most_common(8))}")
    out.append(f"HUMAN: {dict(h['linespacing'].most_common(8))}")
    out.append("\n## Autofit 설정 분포 (텍스트 오버플로 제어)")
    out.append(f"AI   : {dict(a['autofit'])}")
    out.append(f"HUMAN: {dict(h['autofit'])}")

    out.append("\n## 제목 placeholder 좌표 (x,y,w,h cm) — 샘플3")
    out.append(f"AI   : {a['title_coords'][:3]}")
    out.append(f"HUMAN: {h['title_coords'][:3]}")
    out.append("## 서브헤드 placeholder 좌표 — 샘플3")
    out.append(f"AI   : {a['sub_coords'][:3]}")
    out.append(f"HUMAN: {h['sub_coords'][:3]}")

    out.append(f"\n## ★ AI Draft 텍스트 오버플로 추정 ({len(a['overflow'])}건)")
    out.append("(slide, shape, 필요pt, 박스pt, 텍스트)")
    for o in a["overflow"]:
        out.append(f"  {o}")

    rep = "\n".join(out)
    print(rep)
    open(os.path.join(r"C:\Users\cgpar\ppt-skill", "PARADISE_COMPARE.md"), "w", encoding="utf-8").write(rep)
    json.dump({"AI": {k: (dict(v) if isinstance(v, collections.Counter) else v)
                      for k, v in a.items()},
               "HUMAN": {k: (dict(v) if isinstance(v, collections.Counter) else v)
                         for k, v in h.items()}},
              open(os.path.join(r"C:\Users\cgpar\ppt-skill", "paradise_compare.json"), "w", encoding="utf-8"),
              ensure_ascii=False, indent=1)
    print("\nwrote PARADISE_COMPARE.md / paradise_compare.json")


if __name__ == "__main__":
    main()
