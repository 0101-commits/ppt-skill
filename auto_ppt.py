# -*- coding: utf-8 -*-
"""
HCG 제안서 자동 PPT 생성기 v3.0
Design: 실제 final PPTX를 템플릿으로 사용 — 테마/폰트/레이아웃 100% 상속
Content: 롯데알미늄 직무기반 HR제도 설계 및 도입 제안서

핵심 변경 (v3.0):
  - 실제 PPTX 템플릿 기반 → 정확한 색상/폰트 자동 상속
  - Title placeholder → 자동 dark red (#921F0B 계열) 12pt bold
  - Subtitle placeholder (idx=10) → 자동 15pt 맑은 고딕 Semilight
  - 본문 item = noFill rounded rect, 0.5pt border, 맑은 고딕 11pt
  - 좌표 실측값 반영: items y=1.944~, h=0.496", w=4.528", gap=0.589"

실행: python auto_ppt.py
"""

import os, sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ── 경로 ──────────────────────────────────────────────────
BASE = ("C:\\Users\\cgpar\\OneDrive - 휴먼컨설팅그룹\\"
        "09 Admin\\09 etc\\other\\Claude\\롯데알미늄 제안서")
REAL = os.path.join(BASE,
       "[HCG] 롯데알미늄_직무기반 HR제도 설계 및 도입_제안서_final.pptx")
OUT  = "C:\\Users\\cgpar\\ppt-skill\\HCG_Automated_Draft.pptx"

# ── 추출된 실측 좌표 상수 ─────────────────────────────────
# 본문 slid e (layout 2 '본문')
TITLE_Y    = 0.250   # 제목 y (inch)
TITLE_X    = 0.575
TITLE_W    = 9.729
TITLE_H    = 0.305
SUB_Y      = 0.548   # 서브헤드 y
SUB_X      = 0.575
SUB_W      = 9.728
SUB_H      = 0.361

COL_L_X    = 0.691   # 좌측 컬럼 x
COL_R_X    = 5.615   # 우측 컬럼 x
ITEM_W     = 4.528   # 각 아이템 폭
ITEM_H     = 0.496   # 각 아이템 높이
ITEM_Y0    = 1.944   # 첫 번째 아이템 y 시작
ITEM_DY    = 0.589   # 아이템 간격 (y 피치)
HDR_Y      = 1.606   # 컬럼 헤더(OLE 대체) y

LABEL_X    = 0.369   # 행 라벨 x (추진 내용 등)
LABEL_W    = 1.250

# ── 색상 ──────────────────────────────────────────────────
HCG_RED    = RGBColor(0x92, 0x1F, 0x0B)
DARK_GRAY  = RGBColor(0x40, 0x40, 0x40)
MED_GRAY   = RGBColor(0x91, 0x91, 0x91)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
# v4.1 고급 패턴용 (design 스킬 실측 팔레트)
GRAY_DK    = RGBColor(0x1D, 0x1D, 0x1D)
GRAY_LT    = RGBColor(0xD9, 0xD9, 0xD9)
WINE       = RGBColor(0x79, 0x40, 0x39)
CORAL_DK   = RGBColor(0x40, 0x0A, 0x07)
BLUE       = RGBColor(0x35, 0x6C, 0xB5)


# ════════════════════════════════════════════════════════════
# 헬퍼 함수
# ════════════════════════════════════════════════════════════

def _set_font_xml(run, name="맑은 고딕"):
    """run XML에 Korean 폰트 직접 설정"""
    rPr = run._r.get_or_add_rPr()
    for tag in [qn('a:latin'), qn('a:ea')]:
        el = rPr.find(tag)
        if el is None:
            el = etree.SubElement(rPr, tag)
        el.set('typeface', name)

def _set_no_fill(shape):
    """shape fill → noFill (transparent)"""
    sp_el = shape._element
    spPr = sp_el.find(qn('p:spPr'))
    if spPr is None:
        return
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for tag in ['solidFill', 'gradFill', 'blipFill', 'pattFill', 'grpFill', 'noFill']:
        for el in spPr.findall(f'{{{ns}}}{tag}'):
            spPr.remove(el)
    etree.SubElement(spPr, f'{{{ns}}}noFill')

def _set_line_dark(shape, pt=0.5):
    """선 색상 = dk1(black), 두께 pt"""
    sp_el = shape._element
    spPr = sp_el.find(qn('p:spPr'))
    if spPr is None:
        return
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ln_el = spPr.find(f'{{{ns}}}ln')
    if ln_el is None:
        ln_el = etree.SubElement(spPr, f'{{{ns}}}ln')
    ln_el.set('w', str(int(pt * 12700)))
    # color = dk1
    sF = etree.SubElement(ln_el, f'{{{ns}}}solidFill')
    sc = etree.SubElement(sF, f'{{{ns}}}schemeClr')
    sc.set('val', 'dk1')

def _add_run(para, text, fsize=11, bold=False, italic=False,
             color=None, font_name="맑은 고딕"):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(fsize)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    _set_font_xml(run, font_name)
    return run

def _set_line_spacing(para, pct=112):
    """단락 행간 설정 (pct = 100~120...)"""
    pPr = para._p.get_or_add_pPr()
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    lnSpc = pPr.find(f'{{{ns}}}lnSpc')
    if lnSpc is not None:
        pPr.remove(lnSpc)
    lnSpc = etree.SubElement(pPr, f'{{{ns}}}lnSpc')
    spcPct = etree.SubElement(lnSpc, f'{{{ns}}}spcPct')
    spcPct.set('val', str(int(pct * 1000)))


# ── 슬라이드 삭제 ─────────────────────────────────────────
def delete_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    for sld_el in list(sldIdLst):
        rId = sld_el.get(qn('r:id'))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sld_el)


# ── 본문 아이템 shape (rounded rect) ──────────────────────
def add_item(slide, text, left, top, width=ITEM_W, height=ITEM_H,
             fsize=11, bold=False, italic=False, align=PP_ALIGN.LEFT,
             fill_rgb=None, no_border=False, color=None):
    """
    본문 콘텐츠 아이템 shape 추가
    - rounded rectangle (모서리 둥근 직사각형)
    - noFill (기본), 0.5pt dk1 border
    """
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )

    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        _set_no_fill(shape)

    if no_border:
        # 선만 제거 (fill 유지 — _set_no_fill 호출 시 fill까지 사라지던 버그 수정)
        # Remove any line
        sp_el = shape._element
        spPr = sp_el.find(qn('p:spPr'))
        if spPr is not None:
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            for ln_el in spPr.findall(f'{{{ns}}}ln'):
                spPr.remove(ln_el)
            noLn = etree.SubElement(spPr, f'{{{ns}}}ln')
            nf = etree.SubElement(noLn, f'{{{ns}}}noFill')
    else:
        _set_line_dark(shape, 0.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.07)
    tf.margin_right = Inches(0.04)
    tf.margin_top   = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    p = tf.paragraphs[0]
    p.alignment = align
    _add_run(p, text, fsize, bold, italic, color)
    _set_line_spacing(p, 112)

    return shape


# ── 텍스트 박스 ───────────────────────────────────────────
def add_textbox(slide, text, left, top, width, height,
                fsize=11, bold=False, italic=False, align=PP_ALIGN.LEFT,
                color=None, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    _set_no_fill(txBox)
    sp_el = txBox._element
    spPr = sp_el.find(qn('p:spPr'))
    if spPr is not None:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ln_el = etree.SubElement(spPr, f'{{{ns}}}ln')
        nf = etree.SubElement(ln_el, f'{{{ns}}}noFill')

    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02)

    p = tf.paragraphs[0]
    p.alignment = align
    _add_run(p, text, fsize, bold, italic, color)
    _set_line_spacing(p, 112)
    return txBox


# ── 제목/서브헤드 설정 (placeholder 사용) ─────────────────
def set_title(slide, title_text, subtitle_text=None):
    """
    본문 레이아웃 placeholder[0] → 제목 (자동 dark red 스타일)
    placeholder[10] → 서브헤드 (자동 semilight 스타일)
    """
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = title_text
        elif idx == 10 and subtitle_text:
            ph.text = subtitle_text


# ════════════════════════════════════════════════════════════
# 슬라이드 빌더
# ════════════════════════════════════════════════════════════

def build_cover(prs):
    """Slide 1: 표지 (layout 0)"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Title placeholder (ctrTitle, idx=0)
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = "롯데알미늄 직무기반HR 제도 설계 및 도입"
        elif idx == 1:
            ph.text = "- 제안서 -"

    # Date (별도 textbox — placeholder 아래)
    add_textbox(slide, "2026.03",
                1.137, 5.319, 8.563, 0.30,
                fsize=10, bold=False, align=PP_ALIGN.CENTER, color=DARK_GRAY)


def build_toc(prs):
    """Slide 2: 목차 (layout 1)"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # 목차 레이아웃 placeholder 확인
    ph_idxs = [ph.placeholder_format.idx for ph in slide.placeholders]

    # 목차 내용 — 번호 + 섹션명 + 슬라이드 범위
    toc_items = [
        ("I",   "Project Overview / 현황 진단",     "03-04"),
        ("II",  "직무체계 수립",                     "05-10"),
        ("III", "평가제도 개선",                     "11-15"),
        ("IV",  "보상제도 개선",                     "16-19"),
        ("V",   "추진계획 및 HCG 역량",             "20-23"),
    ]

    # 목차 타이틀
    add_textbox(slide, "목차 (Contents)",
                0.5, 0.3, 9.8, 0.5,
                fsize=18, bold=True, color=HCG_RED)

    # 구분선
    from pptx.util import Pt as _Pt
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.5), Inches(0.9), Inches(9.8), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = HCG_RED
    sp_el = line._element
    spPr = sp_el.find(qn('p:spPr'))
    if spPr is not None:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ln_el = etree.SubElement(spPr, f'{{{ns}}}ln')
        nf = etree.SubElement(ln_el, f'{{{ns}}}noFill')

    y = 1.1
    for i, (num, title, pg) in enumerate(toc_items):
        # 번호 원형
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(0.5), Inches(y - 0.04), Inches(0.35), Inches(0.35)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = HCG_RED
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        _add_run(cp, num, fsize=9, bold=True, color=WHITE)

        # 섹션 타이틀
        add_textbox(slide, title,
                    1.0, y, 7.5, 0.4,
                    fsize=13, bold=False, color=DARK_GRAY)
        # 페이지 번호
        add_textbox(slide, pg,
                    8.8, y, 1.2, 0.4,
                    fsize=11, bold=False, align=PP_ALIGN.RIGHT, color=MED_GRAY)
        y += 0.85


def build_overview(prs):
    """Slide 3: Project Overview (layout 2)"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_title(slide,
              "Project Overview",
              "롯데알미늄의 직무체계 및 평가/보상제도 운영방안 도출")

    # 추진 내용 행
    add_textbox(slide, "추진 내용",
                LABEL_X, 2.175, LABEL_W, 0.32,
                fsize=13, italic=True, align=PP_ALIGN.RIGHT, color=MED_GRAY)
    add_textbox(slide, "롯데알미늄의 직무분류체계 수립, 직무기술서 작성(AI 활용), 평가/보상제도 설계를 통한 직무기반 HR제도 도입",
                1.836, 2.175, 8.759, 0.32,
                fsize=11, color=DARK_GRAY)

    # 추진 프로세스 행
    add_textbox(slide, "추진 프로세스",
                0.0, 2.952, 1.618, 0.32,
                fsize=13, italic=True, align=PP_ALIGN.RIGHT, color=MED_GRAY)

    steps = [
        ("현황 진단 및\n방향성 수립",  2.133),
        ("직무분류체계\n수립",          4.161),
        ("평가제도\n개선",              6.168),
        ("보상제도\n개선",              8.174),
    ]
    for title, x in steps:
        sh = add_item(slide, title, x, 2.952, width=1.976, height=0.666,
                      fsize=11, align=PP_ALIGN.CENTER,
                      fill_rgb=HCG_RED, color=WHITE)

    # 화살표 연결선 (간단한 텍스트 →)
    for x in [4.1, 6.1, 8.0]:
        add_textbox(slide, "→",
                    x, 3.18, 0.15, 0.25,
                    fsize=11, align=PP_ALIGN.CENTER, color=MED_GRAY)

    # 세부 설명
    sub_details = [
        (2.133, "현황 파악 Quick Review\n방향성 도출"),
        (4.161, "직무분류체계 초안\nAI 기반 직무기술서"),
        (6.168, "성과관리체계 수립\n평가 운영체계 개선"),
        (8.174, "보상 원칙/전략 수립\n보상 재원 Simulation"),
    ]
    for x, txt in sub_details:
        add_textbox(slide, txt, x, 3.632, 1.951, 0.554,
                    fsize=10, color=DARK_GRAY)

    # 추진 방안
    add_textbox(slide, "추진 방안",
                LABEL_X, 4.654, LABEL_W, 0.50,
                fsize=13, italic=True, align=PP_ALIGN.RIGHT, color=MED_GRAY)
    plan_text = ("추진 기간 : 12주\n"
                 "투입 인력 : PM 1명 + 컨설턴트 2명 (직무체계 전문 컨설턴트 투입)\n"
                 "추진 금액 : 225,000 천원 (부가세 별도; 세부 협의에 따라 변동 가능)")
    add_textbox(slide, plan_text,
                1.836, 4.589, 8.759, 0.972,
                fsize=10, color=DARK_GRAY)


def build_pain_point(prs):
    """Slide 4: Pain Point / 현황 진단 (layout 2) — 2-column"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_title(slide,
              "현황 진단 — Pain Point",
              "롯데알미늄 직무기반 HR제도 도입의 주요 Pain Point 및 개선 방향")

    # 컬럼 헤더
    hdr_items = [
        (COL_L_X, "현재 문제점 (As-Is)",   HCG_RED,   WHITE),
        (COL_R_X, "기대 개선사항 (To-Be)", DARK_GRAY, WHITE),
    ]
    for x, txt, fill, tcol in hdr_items:
        sh = add_item(slide, txt, x, HDR_Y, width=ITEM_W, height=0.333,
                      fsize=11, bold=True, align=PP_ALIGN.CENTER,
                      fill_rgb=fill, color=tcol, no_border=True)

    # Pain Point 데이터 (좌: 문제, 우: 개선)
    pairs = [
        ("직무 데이터 부재 — 체계적 직무 분류/기술 없이 관행적 운영",
         "직무분류체계 수립으로 명확한 직무 Scope 정의"),
        ("노후화된 직무 데이터 — 현업 실태 미반영, 활용 불가",
         "AI 활용 직무기술서 작성으로 신속한 현행화"),
        ("평가 기준 불명확 — 직무 특성 미반영, 공정성 낮음",
         "직무기반 성과관리체계 수립으로 평가 신뢰성 제고"),
        ("보상 운영 기준 부재 — 직무 가치 미반영, 형평성 문제",
         "직무가치 기반 보상체계 설계로 내외부 공정성 확보"),
        ("HR제도 활용 미흡 — 채용/배치/육성 연계 부족",
         "직무기반 통합 HR제도 운영체계 구축"),
    ]

    for i, (left_txt, right_txt) in enumerate(pairs):
        y = ITEM_Y0 + i * ITEM_DY
        add_item(slide, left_txt,  COL_L_X, y, fsize=10, color=DARK_GRAY)
        add_item(slide, right_txt, COL_R_X, y, fsize=10, color=DARK_GRAY)

    # 중간 VS 연결 표시
    for i in range(len(pairs)):
        y = ITEM_Y0 + i * ITEM_DY + 0.05
        add_item(slide, "→", 5.22, y, width=0.394, height=0.394,
                 fsize=10, bold=True, align=PP_ALIGN.CENTER,
                 fill_rgb=HCG_RED, color=WHITE, no_border=True)


def build_section_divider(prs, section_num, section_title, sub):
    """Section 구분 슬라이드 (layout 1 '목차')"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # 섹션 번호 원형
    circle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(1.0), Inches(2.8), Inches(0.9), Inches(0.9)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = HCG_RED
    circle.line.fill.background()
    ctf = circle.text_frame
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    _add_run(cp, section_num, fsize=22, bold=True, color=WHITE)

    # 섹션 타이틀
    add_textbox(slide, section_title,
                2.2, 2.7, 8.0, 0.7,
                fsize=24, bold=True, color=HCG_RED)

    # 서브 설명
    add_textbox(slide, sub,
                2.2, 3.5, 8.0, 0.5,
                fsize=13, color=DARK_GRAY)

    # 구분선
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(2.2), Inches(3.42), Inches(8.0), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = MED_GRAY
    sp_el = line._element
    spPr = sp_el.find(qn('p:spPr'))
    if spPr is not None:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ln_el = etree.SubElement(spPr, f'{{{ns}}}ln')
        nf = etree.SubElement(ln_el, f'{{{ns}}}noFill')


def build_body_2col(prs, title, subtitle, header_l, header_r, left_items, right_items):
    """표준 2컬럼 본문 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_title(slide, title, subtitle)

    # 컬럼 헤더
    add_item(slide, header_l, COL_L_X, HDR_Y, width=ITEM_W, height=0.333,
             fsize=11, bold=True, align=PP_ALIGN.CENTER,
             fill_rgb=HCG_RED, color=WHITE, no_border=True)
    add_item(slide, header_r, COL_R_X, HDR_Y, width=ITEM_W, height=0.333,
             fsize=11, bold=True, align=PP_ALIGN.CENTER,
             fill_rgb=DARK_GRAY, color=WHITE, no_border=True)

    rows = max(len(left_items), len(right_items))
    for i in range(rows):
        y = ITEM_Y0 + i * ITEM_DY
        if i < len(left_items):
            add_item(slide, left_items[i],  COL_L_X, y, fsize=10, color=DARK_GRAY)
        if i < len(right_items):
            add_item(slide, right_items[i], COL_R_X, y, fsize=10, color=DARK_GRAY)

    return slide


def build_body_single(prs, title, subtitle, rows):
    """단일 컬럼 본문 슬라이드 (label + content 구조)"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_title(slide, title, subtitle)

    y = ITEM_Y0
    for label, content in rows:
        if label:
            add_textbox(slide, label,
                        LABEL_X, y, LABEL_W, ITEM_H,
                        fsize=12, italic=True, align=PP_ALIGN.RIGHT, color=MED_GRAY)
        add_item(slide, content,
                 1.836, y, width=8.759, height=ITEM_H,
                 fsize=10, color=DARK_GRAY)
        y += ITEM_DY

    return slide


def build_body_process(prs, title, subtitle, steps, desc_rows=None):
    """프로세스 단계 슬라이드 (→ 화살표 구조)"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_title(slide, title, subtitle)

    step_y = 2.5
    n = len(steps)
    box_w = min(2.2, 9.0 / n)
    start_x = 0.5 + (9.5 - n * box_w - (n-1) * 0.25) / 2

    for i, (step_title, step_detail) in enumerate(steps):
        x = start_x + i * (box_w + 0.25)
        sh = add_item(slide, step_title, x, step_y, width=box_w, height=0.6,
                      fsize=11, bold=True, align=PP_ALIGN.CENTER,
                      fill_rgb=HCG_RED, color=WHITE, no_border=True)
        if step_detail:
            add_textbox(slide, step_detail, x, step_y + 0.65, box_w, 0.8,
                        fsize=9, align=PP_ALIGN.CENTER, color=DARK_GRAY)
        if i < n - 1:
            add_textbox(slide, "→",
                        x + box_w + 0.03, step_y + 0.15, 0.2, 0.3,
                        fsize=13, align=PP_ALIGN.CENTER, color=MED_GRAY)

    if desc_rows:
        y = step_y + 1.6
        for label, content in desc_rows:
            if label:
                add_textbox(slide, label,
                            LABEL_X, y, LABEL_W, 0.38,
                            fsize=11, italic=True, align=PP_ALIGN.RIGHT, color=MED_GRAY)
            add_item(slide, content,
                     1.836, y, width=8.759, height=0.4,
                     fsize=10, color=DARK_GRAY)
            y += 0.46

    return slide


def build_end(prs):
    """Slide 24: End of document (layout 3)"""
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    # 레이아웃이 제공하는 디자인 사용
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            try: ph.text = "Thank you"
            except: pass
    return slide


def build_appendix(prs):
    """Slide 25: Appendix"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_title(slide, "Appendix", "직무분류체계 예시 — 대기업 알미늄 제조업 기준")

    rows = [
        ("대분류", "영업/마케팅  |  생산/제조  |  품질/안전  |  연구/개발  |  경영지원  |  구매/물류"),
        ("중분류",
         "영업관리, 마케팅전략, 고객지원 / 판재생산, 박재생산, 주조 / 품질관리, 안전환경 / 연구개발, 공정개발"),
        ("소분류",
         "영업기획, 국내영업, 해외영업, 마케팅조사, 고객불만처리 / 판재제조, 박재제조, 용해주조, 열처리 / ..."),
        ("직무기술서",
         "각 직무별 Key Responsibilities, Required Competency, Required Knowledge/Skill 정의 (AI 초안 기반)"),
    ]

    y = ITEM_Y0
    for label, content in rows:
        add_textbox(slide, label,
                    LABEL_X, y, LABEL_W, 0.45,
                    fsize=12, italic=True, align=PP_ALIGN.RIGHT, color=MED_GRAY)
        add_item(slide, content,
                 1.836, y, width=8.759, height=0.45,
                 fsize=10, color=DARK_GRAY)
        y += 0.55


# ════════════════════════════════════════════════════════════
# 메인 빌드
# ════════════════════════════════════════════════════════════

# ════════════════════════════════════════════════════════════
# 고급 도식 헬퍼 (v4.1 — 롯데알미늄 _final 실측 패턴)
# ════════════════════════════════════════════════════════════

def add_header_bar(slide, text, x=0.691, y=1.595, w=9.451, h=0.333, fill=None):
    """콘텐츠 영역 헤더 바 (GRAY fill + 흰 bold)."""
    return add_item(slide, text, x, y, width=w, height=h, fsize=11, bold=True,
                    align=PP_ALIGN.LEFT, fill_rgb=fill or MED_GRAY, color=WHITE, no_border=True)

def add_label_tag(slide, text, x, y, w=1.02, h=0.5, fill=None):
    """미니 카테고리 라벨 태그."""
    return add_item(slide, text, x, y, width=w, height=h, fsize=9, bold=True,
                    align=PP_ALIGN.CENTER, fill_rgb=fill or GRAY_DK, color=WHITE, no_border=True)

def add_connector_badge(slide, text, x=4.752, y=1.994, w=1.25, h=0.39, fill=None):
    """중앙 커넥터/키워드 배지."""
    return add_item(slide, text, x, y, width=w, height=h, fsize=9, bold=True,
                    align=PP_ALIGN.CENTER, fill_rgb=fill or MED_GRAY, color=WHITE, no_border=True)

def add_vs_badge(slide, x=5.22, y=1.61, w=0.39, h=0.39, fill=None, text="VS"):
    """Legacy vs HCG 'VS' 원형 배지."""
    sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill or HCG_RED
    sh.line.fill.background()
    p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _add_run(p, text, fsize=10, bold=True, color=WHITE)
    return sh

def add_example_badge(slide, x=9.56, y=1.61, w=0.58, h=0.27):
    """'예시적' 우상단 배지 — 가설/예시 데이터 단정 회피."""
    return add_item(slide, "예시적", x, y, width=w, height=h, fsize=8,
                    align=PP_ALIGN.CENTER, fill_rgb=GRAY_LT, color=DARK_GRAY, no_border=True)

def add_insight_quote(slide, text, x=0.691, y=6.79, w=9.451, h=0.30, color=None):
    """큰따옴표 italic 강조 takeaway."""
    if not (text.startswith("“") or text.startswith('"')):
        text = "“" + text + "”"
    return add_textbox(slide, text, x, y, w, h, fsize=11, bold=True, italic=True,
                       align=PP_ALIGN.CENTER, color=color or HCG_RED)


def build_overview_3col(prs, title, subtitle, cols, bar_label=None, example=False):
    """모듈 Overview — Step1/2/3 3컬럼 + AI 콜아웃. cols=[{header,desc,ai,detail}]×3"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_title(slide, title, subtitle)
    if bar_label:
        add_header_bar(slide, bar_label)
    xs = [0.69, 3.90, 7.11]; w = 3.03
    for i, col in enumerate(cols[:3]):
        x = xs[i]
        add_item(slide, "", x, 2.88, width=w, height=4.17, fill_rgb=GRAY_LT, no_border=True)
        add_item(slide, col.get("header", ""), x, 2.10, width=w, height=0.79, fsize=11, bold=True,
                 align=PP_ALIGN.CENTER, fill_rgb=HCG_RED, color=WHITE, no_border=True)
        add_textbox(slide, col.get("desc", ""), x + 0.10, 2.98, w - 0.20, 0.80, fsize=9, color=DARK_GRAY)
        add_item(slide, col.get("ai", ""), x + 0.10, 3.86, width=w - 0.20, height=1.45, fsize=9, bold=True,
                 align=PP_ALIGN.CENTER, fill_rgb=BLUE, color=WHITE, no_border=True)
        add_textbox(slide, col.get("detail", ""), x + 0.10, 5.45, w - 0.20, 1.45, fsize=8.5, color=DARK_GRAY)
    if example:
        add_example_badge(slide)
    return slide


def build_diff_matrix(prs, title, subtitle, rows, bottom_quote=None, example=False):
    """직군별 차별화 매트릭스. rows=[{group,trait,insight,apply}] (최대 4)"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_title(slide, title, subtitle)
    add_textbox(slide, "업무 특성", 2.00, 2.05, 2.07, 0.30, fsize=9, bold=True, align=PP_ALIGN.CENTER, color=MED_GRAY)
    add_textbox(slide, "핵심 방향성", 4.21, 2.05, 2.68, 0.30, fsize=9, bold=True, align=PP_ALIGN.CENTER, color=MED_GRAY)
    add_textbox(slide, "평가 반영방안", 7.57, 2.05, 2.59, 0.30, fsize=9, bold=True, align=PP_ALIGN.CENTER, color=MED_GRAY)
    ys = [2.46, 3.55, 4.63, 5.74]
    for i, r in enumerate(rows[:4]):
        y = ys[i]
        add_item(slide, r.get("group", ""), 0.69, y, width=1.25, height=0.97, fsize=10, bold=True,
                 align=PP_ALIGN.CENTER, fill_rgb=HCG_RED, color=WHITE, no_border=True)
        add_textbox(slide, r.get("trait", ""), 2.00, y, 2.07, 0.97, fsize=8.5, color=DARK_GRAY)
        ins = r.get("insight", "")
        if not ins.startswith("“"):
            ins = "“" + ins + "”"
        add_item(slide, ins, 4.21, y, width=2.68, height=0.97, fsize=9, bold=True,
                 align=PP_ALIGN.CENTER, fill_rgb=GRAY_LT, color=HCG_RED)
        add_textbox(slide, "▶", 6.99, y + 0.30, 0.47, 0.35, fsize=12, align=PP_ALIGN.CENTER, color=MED_GRAY)
        add_textbox(slide, r.get("apply", ""), 7.57, y, 2.59, 0.97, fsize=8.5, color=DARK_GRAY)
    if bottom_quote:
        add_insight_quote(slide, bottom_quote, 0.69, 6.82, 9.45)
    if example:
        add_example_badge(slide)
    return slide


def build_pain_point_categorized(prs, title, subtitle, left_hdr, right_hdr, rows,
                                 summary_left=None, summary_right=None):
    """Pain Point 카테고리화. rows=[(left_label,left_text,mid,right_text)] (최대 6)"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_title(slide, title, subtitle)
    add_header_bar(slide, left_hdr, 0.691, 1.61, 4.528, 0.333)
    add_header_bar(slide, right_hdr, 5.615, 1.61, 4.528, 0.333)
    y0, dy, h = 1.99, 0.585, 0.50
    for i, row in enumerate(rows[:6]):
        llab, ltext, mid, rtext = row
        y = y0 + i * dy
        add_item(slide, ltext, 1.78, y, width=3.44, height=h, fsize=9, color=DARK_GRAY)
        add_label_tag(slide, llab, 0.691, y, 1.02, h, fill=GRAY_DK)
        if mid:
            add_connector_badge(slide, mid, 4.752, y + 0.05, 1.25, 0.39, fill=MED_GRAY)
        add_item(slide, rtext, 5.615, y, width=4.528, height=h, fsize=9, color=DARK_GRAY)
    if summary_left:
        add_item(slide, summary_left, 0.691, 5.62, width=4.528, height=0.92, fsize=10, bold=True,
                 align=PP_ALIGN.CENTER, fill_rgb=GRAY_LT, color=DARK_GRAY, no_border=True)
    if summary_right:
        add_item(slide, summary_right, 5.615, 5.62, width=4.528, height=0.92, fsize=10, bold=True,
                 align=PP_ALIGN.CENTER, fill_rgb=HCG_RED, color=WHITE, no_border=True)
    return slide


def build_approach_vs(prs, title, subtitle, left_title, left_items, right_title, right_items, bottom_quote=None):
    """Legacy vs HCG 'VS' 대비 블록 + 하단 takeaway."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_title(slide, title, subtitle)
    add_item(slide, "", 0.69, 2.20, width=4.52, height=4.10, fill_rgb=GRAY_LT, no_border=True)
    add_textbox(slide, left_title, 0.80, 2.30, 4.30, 0.40, fsize=12, bold=True, color=DARK_GRAY)
    yy = 2.92
    for it in left_items:
        add_textbox(slide, "• " + it, 0.92, yy, 4.08, 0.70, fsize=10, color=DARK_GRAY); yy += 0.74
    add_item(slide, "", 5.65, 2.20, width=4.45, height=4.10, fill_rgb=None, no_border=False)
    add_textbox(slide, right_title, 5.78, 2.30, 4.20, 0.40, fsize=12, bold=True, color=HCG_RED)
    yy = 2.92
    for it in right_items:
        add_textbox(slide, "• " + it, 5.88, yy, 4.05, 0.70, fsize=10, color=DARK_GRAY); yy += 0.74
    add_vs_badge(slide, 5.22, 3.95, 0.39, 0.39, fill=HCG_RED, text="VS")
    if bottom_quote:
        add_insight_quote(slide, bottom_quote, 1.10, 6.60, 8.64)
    return slide


def build():
    print("Loading template...")
    prs = Presentation(REAL)

    print("Deleting existing slides...")
    delete_all_slides(prs)

    print("Building slides...")

    # Slide 01: 표지
    build_cover(prs)
    print("  S01 Cover ✓")

    # Slide 02: 목차
    build_toc(prs)
    print("  S02 TOC ✓")

    # Slide 03: Project Overview
    build_overview(prs)
    print("  S03 Project Overview ✓")

    # Slide 04: Pain Point
    build_pain_point(prs)
    print("  S04 Pain Point ✓")

    # ── Section I: 직무체계 수립 ──────────────────────────
    build_section_divider(prs, "I", "직무체계 수립",
                          "직무분류체계 설계 및 AI 기반 직무기술서 작성")
    print("  S05 Section I ✓")

    # Slide 06: 직무체계 수립 방향
    build_body_2col(prs,
        "직무체계 수립 방향",
        "직무분류체계 기반의 직무기술서 작성 및 활용체계 구축",
        "현행 문제점", "개선 방향",
        ["직무분류 기준 부재 — 직무와 직책이 혼용",
         "직무 Scope 불명확 — 담당업무 모호",
         "직무기술서 미비 — 구두/관행 의존",
         "직무기반 HR 연계 미흡"],
        ["직무군/직무/세부직무 3단계 분류체계 수립",
         "직무별 Key Roles & Responsibilities 명확화",
         "표준화된 직무기술서 양식 및 AI 작성 지원",
         "채용·평가·보상·육성 연계 체계 구축"]
    )
    print("  S06 직무체계 방향 ✓")

    # Slide 07: 직무분류체계 Framework
    build_body_process(prs,
        "직무분류체계 Framework",
        "대기업 알미늄 제조업 직무 특성을 반영한 3단계 분류체계",
        [("직무군\n(Job Family)",   "영업, 생산, 품질,\n연구, 경영지원 등"),
         ("직무\n(Job)",            "각 직무군 내\n5~15개 직무"),
         ("세부직무\n(Sub-Job)",    "직무별 세분화된\n역할 단위"),
         ("직무기술서\n(JD)",       "Key Roles, 역량,\n자격 요건 기술")],
        [("분류 기준", "업무 유사성, 역량 공통성, 조직 구조 반영"),
         ("적용 범위", "전 사원 대상 (임원 제외), 약 4개 직무군 15개 직무 예상")]
    )
    print("  S07 직무분류 Framework ✓")

    # Slide 08: AI 활용 직무기술서
    build_body_single(prs,
        "AI 활용 직무기술서 작성",
        "ChatGPT/Claude 기반 직무기술서 초안 작성 후 현업 검증 및 Fine-tuning",
        [("작성 방법", "AI 프롬프트 설계 → 초안 생성 → 현업 인터뷰/검증 → 확정"),
         ("작성 항목", "직무 목적, Key Responsibilities(5~8개), 역량, 자격요건, KPI 연계"),
         ("활용 도구", "Claude API / ChatGPT 4o 활용, HCG 표준 JD 템플릿 적용"),
         ("품질 관리", "직무 담당자 및 팀장 검토, HR 최종 확인, 3단계 검증 프로세스")]
    )
    print("  S08 AI 직무기술서 ✓")

    # Slide 09: 직무체계 활용 로드맵
    build_body_process(prs,
        "직무체계 활용 로드맵",
        "수립된 직무체계를 HR 전반에 단계적으로 연계·활용",
        [("Phase 1\n(수립)",        "직무분류체계\n직무기술서 완성"),
         ("Phase 2\n(연계)",        "평가·보상 기준\n직무 연계"),
         ("Phase 3\n(고도화)",      "채용·육성·경력\n통합 운영"),
         ("Phase 4\n(정착)",        "HR Data 축적\n주기적 업데이트")],
    )
    print("  S09 직무체계 로드맵 ✓")

    # Slide 10: 직무기술서 구성 체계
    build_body_single(prs,
        "직무기술서 구성 체계",
        "HCG 표준 JD 양식 기반 — 5개 핵심 구성항목",
        [("직무 개요",  "직무명, 직무군, 소속 부서, 직무 목적 요약 (2~3 문장)"),
         ("핵심 역할",  "Key Responsibilities 5~8개, 각 역할별 수행 활동 기술"),
         ("요구 역량",  "Technical Competency + Behavioral Competency 각 3~5개"),
         ("자격 요건",  "학력, 경력 연수, 필수 자격증, 선호 전공 등"),
         ("KPI 연계",   "직무와 연계된 성과지표 3~5개, 측정 방법 명시")]
    )
    print("  S10 직무기술서 구성 ✓")

    # ── Section II: 평가제도 개선 ─────────────────────────
    build_section_divider(prs, "II", "평가제도 개선",
                          "직무기반 성과관리체계 수립 및 평가 운영체계 개선")
    print("  S11 Section II ✓")

    # Slide 12: 평가제도 현황 진단
    build_body_2col(prs,
        "평가제도 현황 진단",
        "현행 평가제도의 구조적 문제점과 개선 필요성",
        "현행 문제점", "핵심 이슈",
        ["단일 평가지표 — 직무 특성 미반영",
         "평가자 주관성 과다 — 공정성 저하",
         "목표설정 형식화 — KPI 연계 미흡",
         "피드백 체계 부재 — 성장 기회 저해"],
        ["직무 특성에 맞는 차별화된 평가 기준 필요",
         "객관적 평가지표(정량/정성) 균형 필요",
         "목표 → 활동 → 결과의 일관성 확보 필요",
         "주기적 코칭·피드백 문화 정착 필요"]
    )
    print("  S12 평가 현황 ✓")

    # Slide 13: 직무기반 평가모형
    build_body_process(prs,
        "직무기반 평가모형 설계",
        "직무 특성을 반영한 성과(MBO) + 역량(Competency) 이원 평가모형",
        [("목표 설정\n(연초)",   "직무 KPI 3~5개\n목표 수준 합의"),
         ("중간 점검\n(반기)",   "실적 Review\n코칭·피드백"),
         ("최종 평가\n(연말)",   "MBO 60%\n역량 40%"),
         ("결과 활용\n(차년)",   "보상 연계\n육성 계획")]
    )
    print("  S13 평가모형 ✓")

    # Slide 14: 평가지표 체계
    build_body_single(prs,
        "평가지표 체계",
        "직무군별 차별화된 KPI 및 역량지표 설계",
        [("KPI 구조",    "정량 KPI (60%): 매출목표, 생산량, 불량률, 개발건수 등 직무별 차별화"),
         ("역량 지표",   "Behavioral 역량 (40%): 문제해결, 협업, 전문성, 리더십 등 공통+직무"),
         ("등급 체계",   "5등급: S/A/B/C/D → 상위 20% / 중위 60% / 하위 20% 강제배분"),
         ("캘리브레이션", "팀장→부문장 단계적 검토, HR 최종 조정으로 조직 간 형평성 확보")]
    )
    print("  S14 평가지표 ✓")

    # Slide 15: 평가 운영 프로세스
    build_body_process(prs,
        "평가 운영 프로세스",
        "연간 성과관리 Cycle — 목표설정~보상 연계 일원화",
        [("1월\n목표설정",    "개인 KPI\n조직 목표 연계"),
         ("4월\n중간진단",    "1/4분기 실적\n코칭면담"),
         ("7월\n반기평가",    "반기 실적 Review\n하반기 계획"),
         ("12월\n연말평가",   "최종 성과/역량\n종합 평가"),
         ("1월\n결과활용",    "보상 연계\n육성 계획")]
    )
    print("  S15 평가 운영 ✓")

    # ── Section III: 보상제도 개선 ───────────────────────
    build_section_divider(prs, "III", "보상제도 개선",
                          "직무가치 기반 보상체계 설계 및 보상 운영방안 수립")
    print("  S16 Section III ✓")

    # Slide 17: 보상 원칙 및 전략
    build_body_2col(prs,
        "보상 원칙 및 전략",
        "직무가치와 성과를 반영한 보상체계로의 전환",
        "현행 보상 구조", "개선 방향",
        ["연공서열 중심 — 직무가치 미반영",
         "내부 공정성 낮음 — 직무간 보상 불균형",
         "성과 연계 약함 — 성과급 효과 미흡",
         "외부 경쟁력 미검증 — 시장 대비 불명확"],
        ["직무 가치 반영 — Job Grading 기반 Pay Band",
         "내부 공정성 확보 — 직무가치 기반 급여 대역",
         "성과 연계 강화 — 평가 등급별 차별화 보상",
         "외부 경쟁력 확보 — 시장 벤치마킹 주기적 실시"]
    )
    print("  S17 보상 원칙 ✓")

    # Slide 18: 직무기반 보상체계
    build_body_single(prs,
        "직무기반 보상체계 설계",
        "Job Grading → Pay Band → 성과연동 개인 보상 결정",
        [("Job Grading", "직무 복잡성·책임·전문성·영향력 기준 직무 등급(5~7등급) 설계"),
         ("Pay Band",    "등급별 최저-중위-최고 급여 대역 설정 (시장 P50 기준 중위값)"),
         ("성과 연계",   "평가 등급 S/A/B/C/D → 개인 보상 조정 범위 ±15~20% 내 운영"),
         ("인상 원칙",   "시장 인상률 + 성과 + 직무 성장을 반영한 Merit Increase 기준 수립")]
    )
    print("  S18 보상체계 ✓")

    # Slide 19: 보상 시뮬레이션 방향
    build_body_single(prs,
        "보상 시뮬레이션 방향",
        "현행 인건비 총액 유지 원칙 하에 직무가치 기반 급여 재조정",
        [("Simulation 목표",  "현행 인건비 ±5% 이내에서 직무가치 반영 보상 구조 전환"),
         ("시나리오",         "시나리오 A(현행 유지형) / B(점진 전환형) / C(즉시 전환형) 3개 비교"),
         ("조정 기준",        "직무등급 ×Pay Band 중위값 대비 현행 급여 수준 분포 분석"),
         ("협의 사항",        "노조/직원 의견 수렴, 이해관계자 커뮤니케이션 계획 포함")]
    )
    print("  S19 보상 시뮬레이션 ✓")

    # ── Section IV: 추진계획 ──────────────────────────────
    build_section_divider(prs, "IV", "추진계획 및 HCG 역량",
                          "12주 추진 일정, 투입 인력, 제안 비용 및 HCG 수행 역량")
    print("  S20 Section IV ✓")

    # Slide 21: 추진 일정
    build_body_process(prs,
        "추진 일정 (12주)",
        "Phase 1(현황진단 4주) → Phase 2(체계설계 5주) → Phase 3(완료 3주)",
        [("W1-2\n킥오프",     "현황자료 수집\nAS-IS 분석"),
         ("W3-4\n현황진단",   "인터뷰 실시\nPain Point 도출"),
         ("W5-7\n직무설계",   "분류체계 수립\nJD 작성"),
         ("W8-10\nHR제도",    "평가/보상제도\n설계"),
         ("W11-12\n완료",     "최종 보고\nChange Mgmt")]
    )
    print("  S21 추진 일정 ✓")

    # Slide 22: HCG 수행 역량
    build_body_2col(prs,
        "HCG 수행 역량",
        "직무기반 HR제도 설계 전문 역량 보유 — 다수 대기업 구축 경험",
        "수행 분야", "주요 실적",
        ["직무체계 수립",
         "평가제도 개선",
         "보상체계 설계",
         "HR 디지털 전환"],
        ["대기업 제조업 10개사 직무분류체계 구축",
         "평가모형 설계 및 캘리브레이션 운영 지원 15개사",
         "직무기반 Pay Band 설계 8개사",
         "AI 활용 JD 자동화 솔루션 3개사 구축"]
    )
    print("  S22 HCG 역량 ✓")

    # Slide 23: 제안 비용
    build_body_single(prs,
        "제안 비용",
        "롯데알미늄 직무기반 HR제도 설계 및 도입 컨설팅 투입 비용",
        [("총 금액",    "225,000 천원 (VAT 별도) — 세부 범위/투입 형태 협의에 따라 변동 가능"),
         ("인력 구성",  "PM 1명 + Senior 컨설턴트 1명 + 컨설턴트 1명 (총 3명)"),
         ("투입 기간",  "12주 (3개월) — 킥오프~최종 보고"),
         ("산출물",     "직무분류체계(안), 직무기술서(전직무), 평가/보상제도(안), 운영 가이드라인")]
    )
    print("  S23 제안 비용 ✓")

    # Slide 24: End of document
    build_end(prs)
    print("  S24 End ✓")

    # Slide 25: Appendix
    build_appendix(prs)
    print("  S25 Appendix ✓")

    print(f"\nSaving → {OUT}")
    prs.save(OUT)
    print("Done!")
    return OUT


if __name__ == "__main__":
    result = build()
    print(f"\nOutput: {result}")
